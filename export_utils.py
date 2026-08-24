import io
import base64
from datetime import datetime
from typing import List

# ── QR Code ───────────────────────────────────────────────────────────────────
def generate_qr_code(url: str) -> str:
    """Returns base64 PNG string of a QR code for the given URL."""
    import qrcode
    from PIL import Image

    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=10,
        border=4,
    )
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="#0d9488", back_color="white")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return base64.b64encode(buf.getvalue()).decode()


# ── Excel export ──────────────────────────────────────────────────────────────
def export_to_excel(survey, responses, questions) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    # ── Sheet 1: Raw responses ────────────────────────────────────────────────
    ws = wb.active
    ws.title = "Responses"

    teal  = "FF0D9488"
    white = "FFFFFFFF"
    light = "FFf0fdfa"

    header_font  = Font(bold=True, color=white, size=11)
    header_fill  = PatternFill("solid", fgColor=teal)
    center       = Alignment(horizontal="center", vertical="center")
    thin_border  = Border(
        left=Side(style="thin", color="FFE2E8F0"),
        right=Side(style="thin", color="FFE2E8F0"),
        top=Side(style="thin", color="FFE2E8F0"),
        bottom=Side(style="thin", color="FFE2E8F0"),
    )

    # Headers
    base_headers = ["#", "Respondent Name", "Email", "Submitted At", "Completion %"]
    q_headers    = [f"Q{i+1}: {q.text[:40]}" for i, q in enumerate(questions)]
    all_headers  = base_headers + q_headers

    for col_idx, h in enumerate(all_headers, 1):
        cell = ws.cell(row=1, column=col_idx, value=h)
        cell.font      = header_font
        cell.fill      = header_fill
        cell.alignment = center
        cell.border    = thin_border

    ws.row_dimensions[1].height = 28

    # Data rows
    for row_idx, resp in enumerate(responses, 2):
        answer_map = {}
        for ans in resp.answers:
            if ans.values:
                answer_map[ans.question_id] = ", ".join(ans.values)
            else:
                answer_map[ans.question_id] = ans.value or ""

        row_fill = PatternFill("solid", fgColor="FFF8FAFC") if row_idx % 2 == 0 else PatternFill("solid", fgColor=white)
        row_data = [
            row_idx - 1,
            resp.respondent_name or "Anonymous",
            resp.respondent_email or "—",
            resp.submitted_at.strftime("%Y-%m-%d %H:%M"),
            f"{resp.completion_pct:.0f}%",
        ] + [answer_map.get(q.id, "") for q in questions]

        for col_idx, val in enumerate(row_data, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=val)
            cell.fill   = row_fill
            cell.border = thin_border
            cell.alignment = Alignment(vertical="center")

        ws.row_dimensions[row_idx].height = 22

    # Column widths
    for col_idx in range(1, len(all_headers) + 1):
        ws.column_dimensions[get_column_letter(col_idx)].width = 22

    # ── Sheet 2: Summary stats ────────────────────────────────────────────────
    ws2 = wb.create_sheet("Summary")
    ws2["A1"] = f"Survey: {survey.title}"
    ws2["A1"].font = Font(bold=True, size=14, color=teal.lstrip("FF"))
    ws2["A2"] = f"Total Responses: {len(responses)}"
    ws2["A3"] = f"Exported: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── PowerPoint export ─────────────────────────────────────────────────────────
def export_to_pptx(survey, responses, questions) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    TEAL  = RGBColor(0x0D, 0x94, 0x88)
    DARK  = RGBColor(0x0F, 0x17, 0x2A)
    GREY  = RGBColor(0x64, 0x74, 0x8B)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_text_box(slide, text, left, top, width, height, size=18,
                     bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf    = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        return txBox

    def add_rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    sl1 = prs.slides.add_slide(blank)
    add_rect(sl1, 0, 0, 13.33, 7.5, RGBColor(0x0E, 0x11, 0x17))
    add_rect(sl1, 0, 5.8, 13.33, 1.7, TEAL)
    add_text_box(sl1, survey.title, 1, 2, 11, 1.5, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl1, f"{len(responses)} Responses  •  {datetime.utcnow().strftime('%B %d, %Y')}",
                 1, 3.7, 11, 0.8, size=18, color=RGBColor(0xB3, 0xE0, 0xFF), align=PP_ALIGN.CENTER)
    add_text_box(sl1, "Survey Results Report", 1, 6, 11, 0.8, size=14,
                 color=WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 2: Overview ─────────────────────────────────────────────────────
    sl2 = prs.slides.add_slide(blank)
    add_rect(sl2, 0, 0, 13.33, 1.2, TEAL)
    add_text_box(sl2, "Overview", 0.4, 0.2, 12, 0.8, size=28, bold=True, color=WHITE)

    stats = [
        ("📋", "Total Responses", str(len(responses))),
        ("✅", "Completion Rate", f"{sum(r.completion_pct for r in responses)/max(len(responses),1):.0f}%"),
        ("📊", "Total Questions", str(len(questions))),
        ("📬", "Respondents Today", str(sum(1 for r in responses if r.submitted_at.date() == datetime.utcnow().date()))),
    ]

    for i, (icon, label, val) in enumerate(stats):
        x = 0.4 + i * 3.2
        add_rect(sl2, x, 1.5, 2.9, 2.2, RGBColor(0xF0, 0xFD, 0xFA))
        add_text_box(sl2, icon, x + 0.1, 1.6, 2.7, 0.6, size=28)
        add_text_box(sl2, val, x + 0.1, 2.2, 2.7, 0.8, size=30, bold=True, color=TEAL)
        add_text_box(sl2, label, x + 0.1, 3.0, 2.7, 0.5, size=12, color=GREY)

    # ── Slides 3+: One slide per question with answer breakdown ───────────────
    for q_idx, question in enumerate(questions):
        sl = prs.slides.add_slide(blank)
        add_rect(sl, 0, 0, 13.33, 1.2, DARK)
        add_text_box(sl, f"Q{q_idx+1}: {question.text}", 0.4, 0.15, 12.5, 0.9,
                     size=20, bold=True, color=WHITE)

        # Collect answers for this question
        answer_values = []
        for resp in responses:
            for ans in resp.answers:
                if ans.question_id == question.id:
                    if ans.values:
                        answer_values.extend(ans.values)
                    elif ans.value:
                        answer_values.append(ans.value)

        if not answer_values:
            add_text_box(sl, "No responses yet.", 0.4, 1.5, 12, 1, size=16, color=GREY)
            continue

        if question.type in ("multiple_choice", "checkbox", "dropdown", "rating"):
            # Bar chart representation
            from collections import Counter
            counts = Counter(answer_values)
            total  = sum(counts.values())
            items  = sorted(counts.items(), key=lambda x: -x[1])[:8]

            y = 1.4
            max_count = items[0][1] if items else 1
            for label, count in items:
                pct = count / total * 100
                bar_w = (count / max_count) * 9
                add_rect(sl, 2.5, y, bar_w, 0.38, TEAL)
                add_text_box(sl, str(label)[:30], 0.3, y, 2.1, 0.38, size=12, color=DARK)
                add_text_box(sl, f"{count} ({pct:.0f}%)", 2.5 + bar_w + 0.1, y, 2, 0.38, size=12, color=GREY)
                y += 0.52
        else:
            # Text answers — show last 6
            add_text_box(sl, "Sample Responses:", 0.4, 1.4, 12, 0.4, size=13, bold=True, color=TEAL)
            y = 1.9
            for ans_text in answer_values[-6:]:
                add_rect(sl, 0.4, y, 12.5, 0.5, RGBColor(0xF8, 0xFA, 0xFC))
                add_text_box(sl, f"• {str(ans_text)[:120]}", 0.5, y + 0.05, 12.2, 0.4, size=12, color=DARK)
                y += 0.6

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── PDF export ────────────────────────────────────────────────────────────────
def export_to_pdf(survey, responses, questions) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    )
    from collections import Counter

    buf    = io.BytesIO()
    doc    = SimpleDocTemplate(buf, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm,
                               topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    TEAL_C = colors.HexColor("#0d9488")
    DARK_C = colors.HexColor("#0f172a")
    GREY_C = colors.HexColor("#64748b")

    title_style   = ParagraphStyle("Title",   parent=styles["Title"],   textColor=DARK_C, fontSize=22, spaceAfter=6)
    h2_style      = ParagraphStyle("H2",      parent=styles["Heading2"],textColor=TEAL_C, fontSize=14, spaceBefore=16, spaceAfter=6)
    body_style    = ParagraphStyle("Body",    parent=styles["Normal"],  textColor=DARK_C, fontSize=10, spaceAfter=4)
    caption_style = ParagraphStyle("Caption", parent=styles["Normal"],  textColor=GREY_C, fontSize=9)

    story = []

    # Title page
    story.append(Paragraph(survey.title, title_style))
    story.append(Paragraph(f"Survey Results Report  •  {datetime.utcnow().strftime('%B %d, %Y')}", caption_style))
    story.append(HRFlowable(width="100%", thickness=2, color=TEAL_C, spaceAfter=12))

    # Overview table
    avg_completion = sum(r.completion_pct for r in responses) / max(len(responses), 1)
    today_count    = sum(1 for r in responses if r.submitted_at.date() == datetime.utcnow().date())
    overview_data  = [
        ["Metric", "Value"],
        ["Total Responses",   str(len(responses))],
        ["Avg Completion",    f"{avg_completion:.0f}%"],
        ["Total Questions",   str(len(questions))],
        ["Responses Today",   str(today_count)],
    ]
    ov_table = Table(overview_data, colWidths=[8*cm, 6*cm])
    ov_table.setStyle(TableStyle([
        ("BACKGROUND",  (0,0), (-1,0),  TEAL_C),
        ("TEXTCOLOR",   (0,0), (-1,0),  colors.white),
        ("FONTNAME",    (0,0), (-1,0),  "Helvetica-Bold"),
        ("FONTSIZE",    (0,0), (-1,0),  11),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#f0fdfa"), colors.white]),
        ("GRID",        (0,0), (-1,-1), 0.5, colors.HexColor("#e2e8f0")),
        ("PADDING",     (0,0), (-1,-1), 6),
    ]))
    story.append(ov_table)
    story.append(Spacer(1, 0.5*cm))

    # Per-question breakdown
    for q_idx, question in enumerate(questions):
        story.append(Paragraph(f"Q{q_idx+1}: {question.text}", h2_style))

        answer_values = []
        for resp in responses:
            for ans in resp.answers:
                if ans.question_id == question.id:
                    if ans.values:
                        answer_values.extend(ans.values)
                    elif ans.value:
                        answer_values.append(ans.value)

        if not answer_values:
            story.append(Paragraph("No responses yet.", caption_style))
            continue

        if question.type in ("multiple_choice", "checkbox", "dropdown", "rating"):
            from collections import Counter
            counts = Counter(answer_values)
            total  = sum(counts.values())
            tbl_data = [["Option", "Count", "Percentage"]]
            for label, count in sorted(counts.items(), key=lambda x: -x[1]):
                tbl_data.append([str(label), str(count), f"{count/total*100:.1f}%"])

            tbl = Table(tbl_data, colWidths=[9*cm, 3*cm, 4*cm])
            tbl.setStyle(TableStyle([
                ("BACKGROUND",  (0,0), (-1,0),  TEAL_C),
                ("TEXTCOLOR",   (0,0), (-1,0),  colors.white),
                ("FONTNAME",    (0,0), (-1,0),  "Helvetica-Bold"),
                ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#f8fafc"), colors.white]),
                ("GRID",        (0,0), (-1,-1), 0.5, colors.HexColor("#e2e8f0")),
                ("PADDING",     (0,0), (-1,-1), 6),
            ]))
            story.append(tbl)
        else:
            for ans_text in answer_values[-5:]:
                story.append(Paragraph(f"• {str(ans_text)[:200]}", body_style))

        story.append(Spacer(1, 0.3*cm))

    # Respondents table
    if responses:
        story.append(Paragraph("All Respondents", h2_style))
        r_data = [["#", "Name", "Email", "Submitted", "Completion"]]
        for i, resp in enumerate(responses, 1):
            r_data.append([
                str(i),
                resp.respondent_name or "Anonymous",
                resp.respondent_email or "—",
                resp.submitted_at.strftime("%Y-%m-%d %H:%M"),
                f"{resp.completion_pct:.0f}%",
            ])
        r_tbl = Table(r_data, colWidths=[1*cm, 4*cm, 5*cm, 4*cm, 2.5*cm])
        r_tbl.setStyle(TableStyle([
            ("BACKGROUND",  (0,0), (-1,0),  DARK_C),
            ("TEXTCOLOR",   (0,0), (-1,0),  colors.white),
            ("FONTNAME",    (0,0), (-1,0),  "Helvetica-Bold"),
            ("FONTSIZE",    (0,0), (-1,-1), 8),
            ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.HexColor("#f8fafc"), colors.white]),
            ("GRID",        (0,0), (-1,-1), 0.5, colors.HexColor("#e2e8f0")),
            ("PADDING",     (0,0), (-1,-1), 5),
        ]))
        story.append(r_tbl)

    doc.build(story)
    return buf.getvalue()